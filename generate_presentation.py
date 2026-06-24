"""
Génère un PowerPoint de présentation du pipeline RAG agentique Aquila.
Usage : python generate_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Palette de couleurs ─────────────────────────────────────────────────────
C_BG        = RGBColor(0x0F, 0x17, 0x2A)   # bleu nuit (fond)
C_CARD      = RGBColor(0x1A, 0x25, 0x3D)   # bleu marine (carte)
C_ACCENT    = RGBColor(0x3B, 0x82, 0xF6)   # bleu vif (titres/flèches)
C_ACCENT2   = RGBColor(0x10, 0xB9, 0x81)   # vert émeraude (LLM calls)
C_ACCENT3   = RGBColor(0xF5, 0x9E, 0x0B)   # ambre (avertissement/option)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0xAA, 0xB4, 0xC8)
C_DARK_CARD = RGBColor(0x0D, 0x14, 0x22)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # totalement vide
    return prs.slides.add_slide(layout)


def bg(slide, color=C_BG):
    """Remplit le fond de la slide."""
    from pptx.util import Inches
    bg_shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()
    return bg_shape


def box(slide, x, y, w, h, fill=C_CARD, radius=False, border_color=None, border_w=Pt(0)):
    """Ajoute un rectangle (éventuellement avec bordure)."""
    shape_id = 5 if radius else 1  # 5 = rectangle arrondi
    s = slide.shapes.add_shape(shape_id, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = border_w
    else:
        s.line.fill.background()
    return s


def label(slide, text, x, y, w, h,
          font_size=Pt(11), bold=False, color=C_WHITE,
          align=PP_ALIGN.CENTER, wrap=True, italic=False):
    """Ajoute une zone de texte."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def arrow_down(slide, x, y, h=Inches(0.3), color=C_ACCENT):
    """Flèche verticale vers le bas."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        x, y, x, y + h
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


def hline(slide, x, y, w, color=C_ACCENT, width=Pt(1)):
    """Ligne horizontale."""
    c = slide.shapes.add_connector(1, x, y, x + w, y)
    c.line.color.rgb = color
    c.line.width = width


# ════════════════════════════════════════════════════════════════════════════
# Slide 1 — Titre
# ════════════════════════════════════════════════════════════════════════════
def slide_titre(prs):
    sl = blank_slide(prs)
    bg(sl)

    # Bande accent en haut
    top = sl.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.08))
    top.fill.solid(); top.fill.fore_color.rgb = C_ACCENT; top.line.fill.background()

    # Titre principal
    label(sl, "Pipeline RAG Agentique", Inches(1), Inches(1.8), Inches(11.33), Inches(1.2),
          font_size=Pt(44), bold=True, color=C_WHITE)

    # Sous-titre
    label(sl, "Projet Aquila — Architecture complète",
          Inches(1), Inches(3.1), Inches(11.33), Inches(0.6),
          font_size=Pt(22), color=C_ACCENT)

    # Ligne de séparation
    hline(sl, Inches(1), Inches(3.9), Inches(11.33), color=C_ACCENT, width=Pt(1.5))

    # Tags technologies
    tags = ["gemma4:12b", "bge-m3", "ChromaDB", "LangGraph", "BM25 + MMR", "CrossEncoder", "RRF", "HyDE"]
    x_start = Inches(1)
    for i, tag in enumerate(tags):
        col = i % 4
        row = i // 4
        tx = x_start + col * Inches(2.9)
        ty = Inches(4.3) + row * Inches(0.55)
        b = box(sl, tx, ty, Inches(2.5), Inches(0.4), fill=C_CARD,
                border_color=C_ACCENT, border_w=Pt(1))
        label(sl, tag, tx, ty, Inches(2.5), Inches(0.4),
              font_size=Pt(13), color=C_ACCENT, bold=True)

    # Note bas
    label(sl, "4 à 6 appels LLM par question  •  Retrieval hybride  •  Re-ranking CrossEncoder",
          Inches(1), Inches(6.5), Inches(11.33), Inches(0.5),
          font_size=Pt(12), color=C_GRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# Slide 2 — Vue d'ensemble des 3 phases
# ════════════════════════════════════════════════════════════════════════════
def slide_overview(prs):
    sl = blank_slide(prs)
    bg(sl)

    label(sl, "Vue d'ensemble — 3 phases", Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
          font_size=Pt(26), bold=True, color=C_WHITE)
    hline(sl, Inches(0.4), Inches(0.75), Inches(12.5), color=C_ACCENT, width=Pt(1.5))

    phases = [
        {
            "num": "01",
            "titre": "INGESTION",
            "color": C_ACCENT,
            "sub": "Une seule fois (ou après ajout de documents)",
            "steps": [
                "PDF → Markdown  (pymupdf4llm)",
                "Filtre tables des matières",
                "Découpage par titres  (MarkdownHeaderTextSplitter)",
                "Fusion micro-chunks  (< 400 chars)",
                "Découpage par taille  (1 000 chars, overlap 200)",
                "Contextual retrieval  (préfixe LLM + métadonnées)",
                "Vectorisation  (bge-m3 → 1 024 dimensions)",
                "Stockage  (ChromaDB → C:/vector_db_aquila)",
            ]
        },
        {
            "num": "02",
            "titre": "QUESTION / RÉPONSE",
            "color": C_ACCENT2,
            "sub": "À chaque question posée — 4 à 6 appels LLM",
            "steps": [
                "Identification de la source  (LLM #1)",
                "HyDE — réponse fictive  (LLM #2)",
                "Recherche sémantique MMR  (ChromaDB)",
                "Recherche lexicale BM25",
                "Fusion RRF + déduplication Jaccard",
                "Re-ranking CrossEncoder  (top 5)",
                "Évaluation chunks  (LLM #3)",
                "Génération réponse  (LLM #4)",
            ]
        },
        {
            "num": "03",
            "titre": "ÉVALUATION",
            "color": C_ACCENT3,
            "sub": "À la demande — 50 questions de référence",
            "steps": [
                "run_agentic_all.py  → agentic_results.json",
                "evaluate_ragas.py  → ragas_evaluation.csv",
                "Faithfulness  (hallucinations ?)",
                "AnswerRelevancy  (réponse hors sujet ?)",
                "ContextPrecision  (chunks hors sujet ?)",
                "ContextRecall  (info manquante ?)",
                "AnswerCorrectness  (vs. ground truth)",
            ]
        },
    ]

    col_w = Inches(4.1)
    gap   = Inches(0.15)
    x0    = Inches(0.4)

    for i, ph in enumerate(phases):
        cx = x0 + i * (col_w + gap)
        cy = Inches(0.9)
        ch = Inches(6.3)

        # Carte
        box(sl, cx, cy, col_w, ch, fill=C_CARD,
            border_color=ph["color"], border_w=Pt(1.5))

        # En-tête coloré
        header = sl.shapes.add_shape(1, cx, cy, col_w, Inches(0.75))
        header.fill.solid(); header.fill.fore_color.rgb = ph["color"]
        header.line.fill.background()

        label(sl, f"{ph['num']}  {ph['titre']}", cx, cy, col_w, Inches(0.75),
              font_size=Pt(15), bold=True, color=C_BG)
        label(sl, ph["sub"], cx, cy + Inches(0.78), col_w, Inches(0.42),
              font_size=Pt(9.5), color=ph["color"], italic=True)

        for j, step in enumerate(ph["steps"]):
            sy = cy + Inches(1.25) + j * Inches(0.58)
            dot = sl.shapes.add_shape(1, cx + Inches(0.18), sy + Inches(0.1),
                                       Inches(0.08), Inches(0.08))
            dot.fill.solid(); dot.fill.fore_color.rgb = ph["color"]
            dot.line.fill.background()
            label(sl, step, cx + Inches(0.32), sy, col_w - Inches(0.4), Inches(0.52),
                  font_size=Pt(10), color=C_WHITE, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — Phase 1 : Ingestion
# ════════════════════════════════════════════════════════════════════════════
def slide_ingestion(prs):
    sl = blank_slide(prs)
    bg(sl)

    label(sl, "Phase 1 — Ingestion", Inches(0.4), Inches(0.12), Inches(10), Inches(0.55),
          font_size=Pt(26), bold=True, color=C_WHITE)
    label(sl, "Exécutée une seule fois — produit la base vectorielle ChromaDB",
          Inches(0.4), Inches(0.65), Inches(10), Inches(0.35),
          font_size=Pt(12), color=C_GRAY, italic=True, align=PP_ALIGN.LEFT)
    hline(sl, Inches(0.4), Inches(1.0), Inches(12.5), color=C_ACCENT, width=Pt(1.5))

    steps = [
        ("PDF / TXT / DOCX\n(dossier documents/)",
         C_ACCENT, "Source", None),
        ("Extraction Markdown\npymupdf4llm — page par page\nftfy : réparation encodages",
         C_CARD, "Extraction", "Chevauchement inter-pages : 300 chars copiés d'une page à l'autre"),
        ("Filtre tables des matières\n_is_toc_page  (ratio '...' > 30 %)\n→ pages TdM ignorées",
         C_CARD, "Filtre TdM", None),
        ("Découpage par titres\nMarkdownHeaderTextSplitter\n→ métadonnées h1 / h2 / h3",
         C_CARD, "Split titres", "Titres conservés dans le texte (strip_headers=False)"),
        ("Fusion micro-chunks\n_merge_small_chunks\n→ seuil : 400 caractères",
         C_CARD, "Merge", None),
        ("Découpage par taille\nRecursiveCharacterTextSplitter\nchunk=1 000  overlap=200  \\n| protège tableaux",
         C_CARD, "Split taille", None),
        ("Contextual retrieval\n_contextualize_chunks — gemma4:12b\n[source | page | h1>h2>h3 | mots-clés]",
         C_CARD, "Contexte", "1 appel LLM / chunk lors de l'ingestion (checkpoint pickle)"),
        ("Vectorisation + stockage\nbge-m3 → 1 024 dims, lots de 50\n→ C:/vector_db_aquila",
         C_ACCENT, "ChromaDB", None),
    ]

    box_w = Inches(1.45)
    box_h = Inches(0.9)
    gap_x = Inches(0.25)
    y_row1 = Inches(1.15)
    y_row2 = Inches(2.6)
    x0 = Inches(0.3)

    # Ligne 1 : étapes 0-3
    for i in range(4):
        cx = x0 + i * (box_w + gap_x)
        b_color = steps[i][1]
        box(sl, cx, y_row1, box_w, box_h, fill=b_color,
            border_color=C_ACCENT if b_color == C_CARD else None, border_w=Pt(1))
        label(sl, steps[i][0], cx + Inches(0.05), y_row1 + Inches(0.05),
              box_w - Inches(0.1), box_h - Inches(0.1),
              font_size=Pt(8.5), color=C_WHITE if b_color != C_ACCENT else C_BG,
              align=PP_ALIGN.CENTER, bold=(b_color == C_ACCENT))

        if steps[i][3]:
            label(sl, f"↳ {steps[i][3]}", cx, y_row1 + box_h + Inches(0.05),
                  box_w, Inches(0.35),
                  font_size=Pt(7.5), color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

        if i < 3:
            arrow_down(sl, cx + box_w + gap_x / 2, y_row1 + box_h / 3,
                       h=Inches(0.0))
            c = sl.shapes.add_connector(1,
                cx + box_w, y_row1 + box_h / 2,
                cx + box_w + gap_x, y_row1 + box_h / 2)
            c.line.color.rgb = C_ACCENT; c.line.width = Pt(2)

    # Flèche vers la ligne 2
    mid_x = x0 + 3 * (box_w + gap_x) + box_w / 2
    c = sl.shapes.add_connector(1, mid_x, y_row1 + box_h, mid_x, y_row2 - Inches(0.05))
    c.line.color.rgb = C_ACCENT; c.line.width = Pt(2)

    # Ligne 2 : étapes 4-7 (droite à gauche visuellement → on les met de gauche à droite)
    for i in range(4):
        idx = 7 - i  # on dessine de droite à gauche pour simuler le "serpent"
        cx = x0 + i * (box_w + gap_x)
        b_color = steps[idx][1]
        box(sl, cx, y_row2, box_w, box_h, fill=b_color,
            border_color=C_ACCENT2 if idx in [6] else (C_ACCENT if b_color == C_ACCENT else C_ACCENT),
            border_w=Pt(1.5 if idx == 6 else 1))
        txt_color = C_BG if b_color == C_ACCENT else (C_WHITE)
        label(sl, steps[idx][0], cx + Inches(0.05), y_row2 + Inches(0.05),
              box_w - Inches(0.1), box_h - Inches(0.1),
              font_size=Pt(8.5), color=txt_color,
              align=PP_ALIGN.CENTER, bold=(b_color == C_ACCENT))
        if steps[idx][3]:
            label(sl, f"↳ {steps[idx][3]}", cx, y_row2 + box_h + Inches(0.05),
                  box_w, Inches(0.35),
                  font_size=Pt(7.5), color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
        if i < 3:
            c = sl.shapes.add_connector(1,
                cx + box_w, y_row2 + box_h / 2,
                cx + box_w + gap_x, y_row2 + box_h / 2)
            c.line.color.rgb = C_ACCENT2; c.line.width = Pt(2)

    # Légende flèche retour
    label(sl, "←  sens du flux (suite des étapes)", Inches(0.3), y_row2 - Inches(0.3),
          Inches(6), Inches(0.3), font_size=Pt(9), color=C_GRAY, italic=True, align=PP_ALIGN.LEFT)

    # Stats bas de page
    stats = [
        ("~700 chunks", "pour 2 brochures"),
        ("15–25 min", "sur CPU (vectorisation)"),
        ("1 appel LLM/chunk", "pour la contextualisation"),
        ("Checkpoint pickle", "reprise sur crash"),
    ]
    y_stat = Inches(5.85)
    sw = Inches(3.0)
    for i, (val, desc) in enumerate(stats):
        sx = Inches(0.3) + i * (sw + Inches(0.1))
        box(sl, sx, y_stat, sw, Inches(1.0), fill=C_DARK_CARD,
            border_color=C_ACCENT, border_w=Pt(1))
        label(sl, val, sx, y_stat + Inches(0.05), sw, Inches(0.45),
              font_size=Pt(18), bold=True, color=C_ACCENT)
        label(sl, desc, sx, y_stat + Inches(0.5), sw, Inches(0.45),
              font_size=Pt(10), color=C_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# Slide 4 — Phase 2 : Pipeline agentique (schéma central)
# ════════════════════════════════════════════════════════════════════════════
def slide_agentic(prs):
    sl = blank_slide(prs)
    bg(sl)

    label(sl, "Phase 2 — Pipeline Agentique", Inches(0.4), Inches(0.1), Inches(10), Inches(0.55),
          font_size=Pt(24), bold=True, color=C_WHITE)
    label(sl, "LangGraph  •  4 nœuds  •  4 à 6 appels LLM par question",
          Inches(0.4), Inches(0.62), Inches(10), Inches(0.32),
          font_size=Pt(11), color=C_GRAY, italic=True, align=PP_ALIGN.LEFT)
    hline(sl, Inches(0.4), Inches(0.95), Inches(12.5), color=C_ACCENT2, width=Pt(1.5))

    # ── Colonne centrale : flux principal ───────────────────────────────────
    cx = Inches(4.3)
    node_w = Inches(4.6)
    node_h = Inches(0.72)
    gap = Inches(0.32)

    nodes = [
        ("Question de l'utilisateur", C_ACCENT, True),
        ("① identify_sources  [LLM #1]\nChoisit quel(s) PDF concernent la question", C_ACCENT2, False),
        ("② retrieve_node\nHyDE [LLM #2]  →  MMR + BM25  →  RRF  →  Dedup  →  Re-rank", C_CARD, False),
        ("③ grade_documents  [LLM #3]\nLes 5 chunks sont-ils suffisants ?", C_ACCENT3, False),
        ("④ generate_node  [LLM #4]\ngemma4:12b  •  5 chunks  •  temperature=0", C_ACCENT2, False),
        ("Réponse finale", C_ACCENT, True),
    ]

    node_y = [Inches(1.1) + i * (node_h + gap) for i in range(len(nodes))]

    for i, (txt, color, is_io) in enumerate(nodes):
        fill = color if is_io else C_CARD
        border = color
        text_color = C_BG if is_io else C_WHITE
        b_w = Pt(2) if not is_io else Pt(0)
        box(sl, cx, node_y[i], node_w, node_h, fill=fill,
            border_color=border if not is_io else None, border_w=b_w)
        label(sl, txt, cx + Inches(0.1), node_y[i] + Inches(0.05),
              node_w - Inches(0.2), node_h - Inches(0.1),
              font_size=Pt(10 if '\n' in txt else 12),
              bold=is_io, color=text_color, align=PP_ALIGN.CENTER)

        # Flèche vers le suivant
        if i < len(nodes) - 1:
            mid = cx + node_w / 2
            c = sl.shapes.add_connector(1,
                mid, node_y[i] + node_h,
                mid, node_y[i + 1] - Inches(0.02))
            c.line.color.rgb = C_ACCENT; c.line.width = Pt(2)

    # ── Branche reformulation (à droite du grade) ────────────────────────
    grade_idx = 3
    gy = node_y[grade_idx]
    rx = cx + node_w + Inches(0.5)
    rw = Inches(3.3)
    rh = Inches(0.65)
    ry = gy + Inches(0.04)

    # Flèche NON sortant du grade
    c = sl.shapes.add_connector(1,
        cx + node_w, gy + node_h / 2,
        rx, gy + node_h / 2)
    c.line.color.rgb = C_ACCENT3; c.line.width = Pt(2)
    label(sl, "NON\n(< MAX_ATTEMPTS)", cx + node_w + Inches(0.05), gy + node_h / 2 - Inches(0.3),
          Inches(0.9), Inches(0.55), font_size=Pt(8), color=C_ACCENT3, bold=True)

    # Box rewrite_query
    box(sl, rx, ry, rw, rh, fill=C_CARD,
        border_color=C_ACCENT3, border_w=Pt(1.5))
    label(sl, "rewrite_query  [LLM #5]\nReformule la requête sur ce qui manque",
          rx + Inches(0.05), ry + Inches(0.05), rw - Inches(0.1), rh - Inches(0.1),
          font_size=Pt(9.5), color=C_ACCENT3)

    # Flèche rewrite → retrieve
    retrieve_y = node_y[2]
    c = sl.shapes.add_connector(1,
        rx + rw / 2, ry + rh,
        rx + rw / 2, retrieve_y + node_h / 2)
    c.line.color.rgb = C_ACCENT3; c.line.width = Pt(2)
    c2 = sl.shapes.add_connector(1,
        rx + rw / 2, retrieve_y + node_h / 2,
        cx + node_w, retrieve_y + node_h / 2)
    c2.line.color.rgb = C_ACCENT3; c2.line.width = Pt(2)
    label(sl, "HyDE [LLM #6]", rx + Inches(0.1), retrieve_y + node_h / 2 - Inches(0.3),
          Inches(1.8), Inches(0.3), font_size=Pt(8), color=C_ACCENT3, italic=True)

    # Flèche OUI
    label(sl, "OUI", cx - Inches(0.9), node_y[grade_idx] + node_h / 2 - Inches(0.15),
          Inches(0.7), Inches(0.3), font_size=Pt(10), bold=True, color=C_ACCENT2)

    # ── Légende LLM calls ─────────────────────────────────────────────────
    label(sl, "Appels LLM", Inches(0.15), Inches(1.05), Inches(3.8), Inches(0.35),
          font_size=Pt(11), bold=True, color=C_ACCENT2, align=PP_ALIGN.LEFT)
    llm_items = [
        ("#1  identify_sources",  "identifie le(s) PDF"),
        ("#2  HyDE",              "génère une réponse fictive"),
        ("#3  grade_documents",   "évalue les chunks"),
        ("#4  generate_node",     "génère la réponse"),
        ("#5  rewrite_query",     "reformule (si NON)"),
        ("#6  HyDE (2ème tour)",  "nouvelle réponse fictive"),
    ]
    for i, (num, desc) in enumerate(llm_items):
        ly = Inches(1.45) + i * Inches(0.47)
        dot = sl.shapes.add_shape(1, Inches(0.2), ly + Inches(0.1), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_ACCENT2 if i < 4 else C_ACCENT3
        dot.line.fill.background()
        label(sl, num, Inches(0.38), ly, Inches(2.1), Inches(0.38),
              font_size=Pt(9.5), bold=True,
              color=C_ACCENT2 if i < 4 else C_ACCENT3, align=PP_ALIGN.LEFT)
        label(sl, desc, Inches(0.38), ly + Inches(0.2), Inches(3.3), Inches(0.28),
              font_size=Pt(8.5), color=C_GRAY, align=PP_ALIGN.LEFT, italic=True)

    label(sl, "#5 et #6 uniquement si chunks insuffisants",
          Inches(0.15), Inches(4.4), Inches(3.8), Inches(0.35),
          font_size=Pt(8), color=C_ACCENT3, italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — Zoom : Retrieval hybride
# ════════════════════════════════════════════════════════════════════════════
def slide_retrieval(prs):
    sl = blank_slide(prs)
    bg(sl)

    label(sl, "Zoom — Retrieval hybride", Inches(0.4), Inches(0.1), Inches(10), Inches(0.55),
          font_size=Pt(24), bold=True, color=C_WHITE)
    label(sl, "Ce qui se passe à l'intérieur du nœud retrieve_node",
          Inches(0.4), Inches(0.62), Inches(10), Inches(0.32),
          font_size=Pt(11), color=C_GRAY, italic=True, align=PP_ALIGN.LEFT)
    hline(sl, Inches(0.4), Inches(0.95), Inches(12.5), color=C_ACCENT, width=Pt(1.5))

    # HyDE en haut au centre
    hyde_w = Inches(5)
    hyde_x = (SLIDE_W - hyde_w) / 2
    hyde_y = Inches(1.1)
    box(sl, hyde_x, hyde_y, hyde_w, Inches(0.72), fill=C_CARD,
        border_color=C_ACCENT2, border_w=Pt(1.5))
    label(sl, "HyDE  [LLM #2]\nGénère une réponse fictive stylistiquement proche des chunks",
          hyde_x + Inches(0.1), hyde_y + Inches(0.05), hyde_w - Inches(0.2), Inches(0.62),
          font_size=Pt(10), color=C_ACCENT2)

    # Deux branches
    left_x  = Inches(1.2)
    right_x = Inches(7.5)
    branch_w = Inches(4.5)
    branch_y = Inches(2.25)
    branch_h = Inches(1.05)

    # Flèches depuis HyDE
    for tx in [left_x + branch_w / 2, right_x + branch_w / 2]:
        c = sl.shapes.add_connector(1,
            SLIDE_W / 2, hyde_y + Inches(0.72),
            tx, branch_y)
        c.line.color.rgb = C_ACCENT; c.line.width = Pt(1.5)

    # Branche gauche : sémantique MMR
    box(sl, left_x, branch_y, branch_w, branch_h, fill=C_CARD,
        border_color=C_ACCENT, border_w=Pt(1.5))
    label(sl, "Recherche SÉMANTIQUE — MMR\nbge-m3 vectorise la réponse fictive\n25 chunks  (parmi 100 candidats)  •  λ=0.5",
          left_x + Inches(0.1), branch_y + Inches(0.05), branch_w - Inches(0.2), branch_h - Inches(0.1),
          font_size=Pt(9.5), color=C_WHITE)

    # Branche droite : BM25
    box(sl, right_x, branch_y, branch_w, branch_h, fill=C_CARD,
        border_color=C_ACCENT, border_w=Pt(1.5))
    label(sl, "Recherche LEXICALE — BM25\nQuestion originale normalisée (sans accents)\n25 chunks avec correspondances exactes",
          right_x + Inches(0.1), branch_y + Inches(0.05), branch_w - Inches(0.2), branch_h - Inches(0.1),
          font_size=Pt(9.5), color=C_WHITE)

    # Fusion RRF
    rrf_y = Inches(3.75)
    rrf_w = Inches(6)
    rrf_x = (SLIDE_W - rrf_w) / 2

    for bx in [left_x + branch_w / 2, right_x + branch_w / 2]:
        c = sl.shapes.add_connector(1, bx, branch_y + branch_h, rrf_x + rrf_w / 2, rrf_y)
        c.line.color.rgb = C_ACCENT; c.line.width = Pt(1.5)

    box(sl, rrf_x, rrf_y, rrf_w, Inches(0.72), fill=C_CARD,
        border_color=C_ACCENT, border_w=Pt(1.5))
    label(sl, "Fusion RRF  (Reciprocal Rank Fusion)\nscore = 1/(60 + rang_sémantique) + 1/(60 + rang_BM25)  →  15 candidats",
          rrf_x + Inches(0.1), rrf_y + Inches(0.05), rrf_w - Inches(0.2), Inches(0.62),
          font_size=Pt(9.5), color=C_WHITE)

    # Étapes suivantes en ligne
    next_steps = [
        ("Déduplication\nJaccard > 80 %\n→ quasi-doublons écartés", C_ACCENT3),
        ("Re-ranking\nCrossEncoder\nBAAI/bge-reranker-v2-m3", C_ACCENT),
        ("Top 5 chunks\npour le LLM\n(seuil logit = 0.0)", C_ACCENT2),
    ]
    step_w = Inches(3.5)
    step_h = Inches(0.95)
    step_y = Inches(4.9)
    step_x0 = (SLIDE_W - (3 * step_w + 2 * Inches(0.3))) / 2

    prev_mid = rrf_x + rrf_w / 2
    c = sl.shapes.add_connector(1, prev_mid, rrf_y + Inches(0.72),
                                  prev_mid, step_y)
    c.line.color.rgb = C_ACCENT; c.line.width = Pt(1.5)

    for i, (txt, color) in enumerate(next_steps):
        sx = step_x0 + i * (step_w + Inches(0.3))
        box(sl, sx, step_y, step_w, step_h, fill=C_CARD,
            border_color=color, border_w=Pt(1.5))
        label(sl, txt, sx + Inches(0.1), step_y + Inches(0.05),
              step_w - Inches(0.2), step_h - Inches(0.1),
              font_size=Pt(9.5), color=C_WHITE)
        if i < 2:
            c = sl.shapes.add_connector(1,
                sx + step_w, step_y + step_h / 2,
                sx + step_w + Inches(0.3), step_y + step_h / 2)
            c.line.color.rgb = color; c.line.width = Pt(2)

    # Notes de bas
    notes = [
        "MMR : Maximal Marginal Relevance — pertinence ET diversité (évite 25 extraits du même paragraphe)",
        "RRF : indépendant des scores bruts, ne regarde que les rangs — robuste aux différences d'échelle MMR vs BM25",
        "CrossEncoder : lit (question + chunk) ensemble — plus précis qu'un embedding mais plus lent",
    ]
    for i, note in enumerate(notes):
        label(sl, f"• {note}", Inches(0.4), Inches(6.05) + i * Inches(0.38),
              Inches(12.5), Inches(0.35),
              font_size=Pt(8.5), color=C_GRAY, align=PP_ALIGN.LEFT, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — Évaluation RAGAS
# ════════════════════════════════════════════════════════════════════════════
def slide_evaluation(prs):
    sl = blank_slide(prs)
    bg(sl)

    label(sl, "Phase 3 — Évaluation RAGAS", Inches(0.4), Inches(0.1), Inches(10), Inches(0.55),
          font_size=Pt(24), bold=True, color=C_WHITE)
    label(sl, "50 questions de référence  •  5 métriques  •  Modèles locaux (Ollama)",
          Inches(0.4), Inches(0.62), Inches(10), Inches(0.32),
          font_size=Pt(11), color=C_GRAY, italic=True, align=PP_ALIGN.LEFT)
    hline(sl, Inches(0.4), Inches(0.95), Inches(12.5), color=C_ACCENT3, width=Pt(1.5))

    # Flux évaluation
    eval_steps = [
        ("questions.json\n50 questions + réponses\nde référence", C_ACCENT3),
        ("run_agentic_all.py\nLance le pipeline sur\nchaque question", C_CARD),
        ("agentic_results.json\nRéponse LLM + chunks\n+ logs par question", C_ACCENT3),
        ("evaluate_ragas.py\nConstruit le dataset\nRAGAS et score", C_CARD),
        ("ragas_evaluation.csv\n5 métriques\npar question", C_ACCENT3),
    ]
    ew = Inches(2.3)
    eh = Inches(1.0)
    ey = Inches(1.2)
    ex0 = Inches(0.4)
    gap = Inches(0.2)

    for i, (txt, color) in enumerate(eval_steps):
        ex = ex0 + i * (ew + gap + Inches(0.3))
        box(sl, ex, ey, ew, eh, fill=color if color == C_ACCENT3 else C_CARD,
            border_color=color, border_w=Pt(1.5))
        t_color = C_BG if color == C_ACCENT3 else C_WHITE
        label(sl, txt, ex + Inches(0.08), ey + Inches(0.05), ew - Inches(0.16), eh - Inches(0.1),
              font_size=Pt(9), color=t_color, bold=(color == C_ACCENT3))
        if i < len(eval_steps) - 1:
            ax = ex + ew
            c = sl.shapes.add_connector(1, ax, ey + eh / 2, ax + gap + Inches(0.3), ey + eh / 2)
            c.line.color.rgb = C_ACCENT3; c.line.width = Pt(2)

    # 5 métriques
    label(sl, "Les 5 métriques RAGAS", Inches(0.4), Inches(2.5), Inches(12), Inches(0.45),
          font_size=Pt(16), bold=True, color=C_ACCENT3, align=PP_ALIGN.LEFT)

    metrics = [
        ("Faithfulness", "Le LLM invente-t-il ?\nScore bas → hallucinations", "Non", C_ACCENT),
        ("AnswerRelevancy", "La réponse répond-elle\nà la question ?", "Non", C_ACCENT),
        ("ContextPrecision", "Les chunks récupérés\nsont-ils pertinents ?", "Non", C_ACCENT),
        ("ContextRecall", "Les chunks couvrent-ils\ntout ce que contient la référence ?", "Oui", C_ACCENT2),
        ("AnswerCorrectness", "La réponse est-elle\ncorrecte vs. référence ?", "Oui", C_ACCENT2),
    ]
    mw = Inches(2.4)
    mh = Inches(2.5)
    mx0 = Inches(0.35)
    my = Inches(3.05)

    for i, (name, desc, gt, color) in enumerate(metrics):
        mx = mx0 + i * (mw + Inches(0.12))
        box(sl, mx, my, mw, mh, fill=C_CARD, border_color=color, border_w=Pt(1.5))

        # Header
        header = sl.shapes.add_shape(1, mx, my, mw, Inches(0.45))
        header.fill.solid(); header.fill.fore_color.rgb = color; header.line.fill.background()
        label(sl, name, mx, my, mw, Inches(0.45),
              font_size=Pt(11), bold=True, color=C_BG)

        label(sl, desc, mx + Inches(0.1), my + Inches(0.5), mw - Inches(0.2), Inches(1.2),
              font_size=Pt(10), color=C_WHITE, align=PP_ALIGN.LEFT)

        gt_color = C_ACCENT2 if gt == "Oui" else C_GRAY
        label(sl, f"Ground truth : {gt}", mx + Inches(0.1), my + Inches(1.8),
              mw - Inches(0.2), Inches(0.4),
              font_size=Pt(9), color=gt_color, italic=True, align=PP_ALIGN.LEFT)

        label(sl, "0.0 → 1.0", mx + Inches(0.1), my + mh - Inches(0.45),
              mw - Inches(0.2), Inches(0.35),
              font_size=Pt(10), bold=True, color=color)

    # Note interprétation
    label(sl, "Scores faibles → Faithfulness : hallucinations  •  ContextPrecision : retrieval hors-sujet  •  ContextRecall : info manquante",
          Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.35),
          font_size=Pt(9), color=C_GRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# Assemblage
# ════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_titre(prs)
    slide_overview(prs)
    slide_ingestion(prs)
    slide_agentic(prs)
    slide_retrieval(prs)
    slide_evaluation(prs)

    out = "documentation/RAG_Agentique_Presentation.pptx"
    prs.save(out)
    print(f"OK - Fichier genere : {out}")


if __name__ == "__main__":
    main()
